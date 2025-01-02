import logging
import os

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from models.presentation_config import PresentationConfig
from src.content_generator import ContentGenerator


class SlideContentHandler:
    def __init__(self, content_generator: ContentGenerator):
        logging.info("Initializing SlideContentHandler")
        self.content_generator = content_generator

        # Convert to PowerPoint inches (1 inch = 914400 EMUs)
        self.SLIDE_WIDTH = Inches(13.333)
        self.SLIDE_HEIGHT = Inches(7.5)

        # Add standard measurements
        self.HEADER_HEIGHT = Inches(1.2)
        self.FOOTER_HEIGHT = Inches(0.4)
        self.MARGIN = Inches(0.5)
        self.CONTENT_TOP = Inches(1.5)
        self.STANDARD_TEXT_SIZE = Pt(16)

    def get_logo_image_path(self, presentation: PresentationConfig):
        logging.info("Getting logo image path")
        if presentation.logo_base64:
            logging.debug("Using base64 logo")
            path = self.content_generator.save_base64_image(
                presentation.logo_base64,
                f"{presentation.output_path}/screenshot/logo.png",
            )
            if not path:
                raise ValueError("Failed to save logo image")
            return path
        else:
            return self.content_generator.generate_image_from_text(
                text=f"logo_description: {presentation.logo_description}",
                output_path=f"{presentation.output_path}/screenshot/logo.png",
                height=200,  # Smaller size for logo
                width=200,
            )

    def add_slide(
        self, prs, layout, slide_data, presentation: PresentationConfig, lock=None
    ):
        logging.info(f"Adding slide: {slide_data.title} (Style: {slide_data.style})")
        if lock:
            with lock:
                slide = prs.slides.add_slide(layout)
        else:
            slide = prs.slides.add_slide(layout)
        if slide_data.style == "cover":
            self._add_cover_slide(slide, slide_data, presentation, prs)
        else:
            self._add_content_slide(slide, slide_data, presentation)
            self._add_footer(slide, presentation)

    def _add_cover_slide(self, slide, slide_data, presentation, prs):
        logging.info("Adding cover slide")

        # Add background shape for the entire slide
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        background.line.fill.background()

        # Calculate dimensions for left and right sections
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        left_section_width = slide_width * 0.55  # Use 55% for image section
        right_section_width = slide_width * 0.45  # Use 45% for text/logo section

        # Add the image on the left section
        cover_image_path = self.content_generator.generate_image_from_text(
            text=f"Title: {slide_data.title} \n Subtitle: {slide_data.subtitle} Based on the the title and subtitle, create a professional presentation cover image",
            output_path=f"{presentation.output_path}/screenshot/cover.png",
        )
        slide.shapes.add_picture(
            cover_image_path,
            Inches(0),  # Start from left edge
            Inches(0),
            width=left_section_width,
            height=slide_height,
        )

        # Add title with adjusted position on right side
        right_margin = left_section_width + Inches(0.5)
        title_width = right_section_width - Inches(1)

        title_box = slide.shapes.add_textbox(
            right_margin, Inches(1.5), title_width, Inches(2.5)
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_run = title_paragraph.add_run()
        title_run.text = slide_data.title
        title_run.font.size = Pt(40)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(0, 0, 0)  # Black text

        # Add subtitle
        if slide_data.subtitle:
            subtitle_box = slide.shapes.add_textbox(
                right_margin, Inches(4), title_width, Inches(1.5)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.clear()
            subtitle_frame.word_wrap = True
            subtitle_paragraph = subtitle_frame.paragraphs[0]
            subtitle_paragraph.alignment = PP_ALIGN.LEFT
            subtitle_run = subtitle_paragraph.add_run()
            subtitle_run.text = slide_data.subtitle
            subtitle_run.font.size = Pt(28)
            subtitle_run.font.color.rgb = RGBColor(0, 0, 0)  # Black text

        # Add logo at bottom right corner if available
        if presentation.logo_path and os.path.exists(presentation.logo_path):
            logo_width = Inches(3)  # Larger logo size
            logo_height = Inches(2)  # Maintain aspect ratio
            logo_left = (
                left_section_width + (right_section_width - logo_width) / 2
            )  # Center in right section
            logo_top = slide_height - logo_height - Inches(1)  # Bottom margin

            slide.shapes.add_picture(
                presentation.logo_path,
                logo_left,
                logo_top,
                width=logo_width,
                height=logo_height,
            )

        # Add comments to slide notes if they exist
        if hasattr(slide_data, "comments") and slide_data.comments:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data.comments

    def _add_content_slide(self, slide, slide_data, presentation: PresentationConfig):
        logging.info(f"Adding content slide: {slide_data.title}")
        logging.debug(f"Generating background image for style: {slide_data.style}")
        # Adjust section widths for better proportions
        left_section_width = self.SLIDE_WIDTH * 0.6  # 50% for content
        right_section_width = self.SLIDE_WIDTH * 0.4  # 50% for image

        # Generate slide background image
        image_prompt = f"Title: {slide_data.title}\nStyle: {slide_data.style}\n Content: {slide_data.content}"
        background_image_path = self.content_generator.generate_image_from_text(
            text=image_prompt,
            output_path=f"{presentation.output_path}/screenshot/slide_{slide_data.title.replace(' ', '_')}.png",
        )

        if background_image_path:
            # Calculate image dimensions
            image_margin = Inches(0.2)
            image_left = left_section_width + image_margin
            image_top = self.CONTENT_TOP
            image_width = right_section_width - (image_margin * 2)
            image_height = self.SLIDE_HEIGHT * 0.7

            # Replace the manual image addition with the add_image helper method
            self.add_image(
                slide,
                background_image_path,
                image_left,
                image_top,
                presentation,
                width=image_width,
                height=image_height,
            )

        # Fix title positioning and ensure it's added
        title_box = slide.shapes.title
        if not title_box:  # If title placeholder doesn't exist, create it
            title_box = slide.shapes.add_textbox(
                self.MARGIN, self.MARGIN, int(left_section_width), self.HEADER_HEIGHT
            )
        else:
            title_box.left = self.MARGIN
            title_box.top = self.MARGIN
            title_box.width = int(left_section_width)

        title_frame = title_box.text_frame
        title_frame.clear()  # Clear existing text
        paragraph = title_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = slide_data.title
        run.font.name = presentation.theme.fonts.title.name
        run.font.size = Pt(presentation.theme.fonts.title.size)
        run.font.color.rgb = RGBColor(
            presentation.theme.colors.title.r,
            presentation.theme.colors.title.g,
            presentation.theme.colors.title.b,
        )

        # Add content below title with reduced spacing
        content_left = self.MARGIN  # Changed from Inches(1) to align with title
        if slide_data.style == "bullets":
            self._add_bullet_content(
                slide, slide_data, presentation, left_section_width, content_left
            )
        elif slide_data.style == "table":
            self._add_table_content(slide, slide_data, presentation, left_section_width)
        elif slide_data.style == "paragraph":
            self._add_paragraph_content(
                slide, slide_data, presentation, left_section_width
            )

        # Add comments to slide notes if they exist
        if hasattr(slide_data, "comments") and slide_data.comments:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data.comments

    def _add_footer(self, slide, presentation):
        logging.info("Adding footer to slide")
        """Add footer with logo to slide"""
        # Set footer dimensions and position
        footer_height = Inches(0.4)
        footer_top = self.SLIDE_HEIGHT - footer_height  # Remove extra margin

        # Add footer background shape
        footer_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), footer_top, self.SLIDE_WIDTH, footer_height
        )
        footer_bg.fill.solid()
        footer_bg.fill.fore_color.rgb = RGBColor(245, 245, 245)
        footer_bg.line.width = Pt(0)  # Remove border

        # Add footer text
        text_margin = Inches(0.5)
        footer_text_width = self.SLIDE_WIDTH - (text_margin * 2)
        if presentation.logo_path and os.path.exists(presentation.logo_path):
            footer_text_width -= Inches(1.5)  # Reserve space for logo

        footer_text = slide.shapes.add_textbox(
            text_margin,
            footer_top,  # Align exactly with bottom
            footer_text_width,
            footer_height,
        )

        text_frame = footer_text.text_frame
        text_frame.word_wrap = True
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = presentation.theme.footer
        run.font.name = presentation.theme.fonts.footer.name
        run.font.size = Pt(presentation.theme.fonts.footer.size)
        run.font.color.rgb = RGBColor(
            presentation.theme.colors.footer.r,
            presentation.theme.colors.footer.g,
            presentation.theme.colors.footer.b,
        )

        # Add logo if available
        if presentation.logo_path and os.path.exists(presentation.logo_path):
            logo_size = Inches(0.3)
            logo_left = self.SLIDE_WIDTH - logo_size - Inches(0.5)
            logo_top = footer_top + Inches(0.05)  # Adjusted vertical position

            slide.shapes.add_picture(
                presentation.logo_path,
                logo_left,
                logo_top,
                width=logo_size,
                height=logo_size,
            )

    def _add_bullet_content(
        self, slide, slide_data, presentation, section_width, content_left
    ):
        logging.info("Adding bullet content")

        content_box = slide.shapes.add_textbox(
            content_left,
            self.CONTENT_TOP - Inches(0.3),  # Reduced space between title and content
            section_width - (self.MARGIN * 2),
            self.SLIDE_HEIGHT - self.CONTENT_TOP - self.FOOTER_HEIGHT,
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        # Remove any empty strings at the start of content
        content = [item for item in slide_data.content if str(item).strip()]

        for item in content:
            p = text_frame.add_paragraph()
            p.level = item.level if hasattr(item, "level") else 0

            # Reduced spacing between bullet points
            if p.level == 0:
                p.space_before = Pt(6)  # Reduced from 12
                p.space_after = Pt(3)  # Reduced from 6
            else:
                p.space_before = Pt(3)  # Reduced from 6
                p.space_after = Pt(3)  # Reduced from 6

            text = item.text if hasattr(item, "text") else str(item)

            if " - " in text:
                header, content = text.split(" - ", 1)
                # Add header (bold)
                run = p.add_run()
                run.text = header.strip()
                run.font.bold = True
                run.font.size = self.STANDARD_TEXT_SIZE
                run.font.name = presentation.theme.fonts.text.name

                # Add separator
                run = p.add_run()
                run.text = " - "
                run.font.size = self.STANDARD_TEXT_SIZE
                run.font.name = presentation.theme.fonts.text.name

                # Add content (normal)
                run = p.add_run()
                run.text = content.strip()
                run.font.size = self.STANDARD_TEXT_SIZE
                run.font.name = presentation.theme.fonts.text.name
            else:
                # If no dash separator, treat entire text as content
                run = p.add_run()
                run.text = text
                run.font.size = self.STANDARD_TEXT_SIZE
                run.font.name = presentation.theme.fonts.text.name

    def _add_table_content(self, slide, slide_data, presentation, left_section_width):
        logging.info("Adding table content")
        logging.debug(
            f"Table dimensions: {len(slide_data.content.rows)+1}x{len(slide_data.content.headers)}"
        )
        rows = len(slide_data.content.rows) + 1  # +1 for header
        cols = len(slide_data.content.headers)

        # Calculate available space using standard measurements
        max_table_height = self.SLIDE_HEIGHT - self.CONTENT_TOP - self.FOOTER_HEIGHT

        table = slide.shapes.add_table(
            rows,
            cols,
            self.MARGIN * 2,
            self.CONTENT_TOP,
            left_section_width - (self.MARGIN * 2),
            min(Inches(0.5 * rows), max_table_height),
        ).table

        # Add headers with standard text size
        for idx, header in enumerate(slide_data.content.headers):
            cell = table.cell(0, idx)
            cell.text = header
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = self.STANDARD_TEXT_SIZE

        # Add data with standard text size
        for row_idx, row in enumerate(slide_data.content.rows, start=1):
            for col_idx, cell_text in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_text)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = self.STANDARD_TEXT_SIZE

    def _add_paragraph_content(
        self, slide, slide_data, presentation, left_section_width
    ):
        logging.info("Adding paragraph content")

        content_box = slide.shapes.add_textbox(
            self.MARGIN * 2,
            self.CONTENT_TOP,
            left_section_width - (self.MARGIN * 2),
            self.SLIDE_HEIGHT - self.CONTENT_TOP - self.FOOTER_HEIGHT,
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.add_paragraph()
        # Remove any leading/trailing whitespace from content
        p.text = str(slide_data.content).strip()
        p.font.size = self.STANDARD_TEXT_SIZE
        p.font.name = presentation.theme.fonts.text.name

    def add_image(
        self, slide, image_path, left, top, presentation, width=None, height=None
    ):
        logging.info(f"Adding image: {image_path}")
        """Add image to slide with proper sizing and rounded corners"""
        if width is None:
            width = Inches(6)  # Default width
        if height is None:
            height = Inches(4)  # Default height

        # Adjust circle position to bottom left
        circle_size = Inches(0.8)
        overlap = circle_size * 0.8  # Increased overlap to hide more of the circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left - (circle_size - (overlap * 1.2)),  # Moved to left side
            top + height - (circle_size - (overlap * 0.4)),  # Moved to bottom
            circle_size,
            circle_size,
        )

        # Match circle style to rectangle
        circle.fill.solid()
        theme_color = presentation.theme.colors.bullet
        circle.fill.fore_color.rgb = RGBColor(
            theme_color.r, theme_color.g, theme_color.b
        )
        circle.line.width = Pt(0)

        # Then add the rounded rectangle (it will be on top of the circle)
        rounded_rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left,
            top,
            width,
            height,
        )

        # Set the fill to be white (or any color that matches your background)
        rounded_rect.fill.solid()
        rounded_rect.fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Set the border properties
        theme_color = presentation.theme.colors.bullet
        rounded_rect.line.color.rgb = RGBColor(
            theme_color.r, theme_color.g, theme_color.b
        )
        rounded_rect.line.width = Pt(2)

        # Calculate padding for the image (10% of the shape size)
        padding_w = width * 0.1
        padding_h = height * 0.1

        # Add image last (it should be on top)
        image = slide.shapes.add_picture(
            image_path,
            left + padding_w,
            top + padding_h,
            width - (padding_w * 2),
            height - (padding_h * 2),
        )

        # Set explicit z-order to ensure proper layering
        circle.z_order = 3  # Move to back
        rounded_rect.z_order = 2  # Middle layer
        image.z_order = 1  # Front layer

        return image
