import logging
import os

import boto3
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches

from src.content_generator import ContentGenerator
from src.serper_searcher import SerperAgent
from src.slide_content_handler import SlideContentHandler
from src.storage_handler import StorageHandler
from src.youtube_searcher import YouTubeAgent

# Configure logging at the top of the file
logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)
logger = logging.getLogger(__name__)


# Define PP_Z_ORDER enum since it's not available in python-pptx
class PP_Z_ORDER:
    """Z-order constants for shapes"""

    BRING_TO_FRONT = 0
    SEND_TO_BACK = 1
    BRING_FORWARD = 2
    SEND_BACKWARD = 3


class PresentationMaker:
    def __init__(self):
        # Initialize AWS connection
        self.output_path = ""

        self.aws_access_key_id = os.getenv("AWS_ACCESS_KEY_ID", "")
        self.aws_secret_access_key = os.getenv("AWS_SECRET_ACCESS_KEY", "")
        self.aws_region = os.getenv("AWS_DEFAULT_REGION", "us-east-1")

        # Update debug logging for AWS credentials
        logger.info(f"AWS Region: {self.aws_region}")
        logger.info(f"Has AWS Access Key: {bool(self.aws_access_key_id)}")
        logger.info(f"Has AWS Secret Key: {bool(self.aws_secret_access_key)}")
        logger.debug(f"Environment variables present: {list(os.environ.keys())}")

        # Initialize AWS session
        if self.aws_access_key_id and self.aws_secret_access_key:
            logger.info("Initializing AWS session with access keys")
            try:
                self.session = boto3.Session(
                    aws_access_key_id=self.aws_access_key_id,
                    aws_secret_access_key=self.aws_secret_access_key,
                    region_name=self.aws_region,
                )
                logger.info("Successfully created AWS session with access keys")
            except Exception as e:
                logger.error(f"Error creating AWS session with access keys: {str(e)}")
                raise
        else:
            self.aws_profile = os.getenv("AWS_PROFILE", "")
            print(f"No access keys found, checking for AWS profile: {self.aws_profile}")
            print(f"Available AWS profiles: {boto3.Session().available_profiles}")

            if self.aws_profile:
                print(f"Initializing AWS session with profile: {self.aws_profile}")
                try:
                    self.session = boto3.Session(
                        profile_name=self.aws_profile, region_name=self.aws_region
                    )
                    print("Successfully created AWS session with profile")
                except Exception as e:
                    print(f"Error creating AWS session with profile: {str(e)}")
                    raise
            else:
                print("Warning: No AWS credentials or profile found")

        # Initialize handlers
        self.storage_handler = StorageHandler()
        self.content_generator = ContentGenerator(self.session)
        self.slide_handler = SlideContentHandler(self.content_generator)

        # Set 16:9 aspect ratio dimensions
        self.SLIDE_WIDTH = Inches(13.333)
        self.SLIDE_HEIGHT = Inches(7.5)

        # Add theme color constants
        self.THEME_COLORS = [
            MSO_THEME_COLOR.ACCENT_1,
            MSO_THEME_COLOR.ACCENT_2,
            MSO_THEME_COLOR.ACCENT_3,
            MSO_THEME_COLOR.ACCENT_4,
            MSO_THEME_COLOR.ACCENT_5,
            MSO_THEME_COLOR.ACCENT_6,
        ]

        # Add slide layout dimensions for split design
        self.LEFT_MARGIN = Inches(1)
        self.TITLE_WIDTH = Inches(4)  # Width of the title section on the left
        self.CONTENT_LEFT = Inches(1)  # Left margin for content
        self.CONTENT_WIDTH = Inches(11)  # Width for content area

    def create_presentation(self, config_path):
        print(f"\nStarting presentation creation with config: {config_path}")
        try:
            # Load configuration
            print("Loading presentation configuration...")
            if not os.path.exists(config_path):
                print(f"ERROR: Configuration file not found at {config_path}")
                raise FileNotFoundError(
                    f"Configuration file not found at {config_path}"
                )

            with open(config_path, "r") as f:
                print(f"Raw config content: {f.read()}")

            presentation = self.storage_handler.load_presentation_config(config_path)
            print(f"Loaded configuration: {presentation.dict()}")

            print(f"Clearing output folder: {presentation.output_path}")
            # self.storage_handler.clear_output_folder(presentation.output_path)

            # Create presentation with 16:9 aspect ratio
            print("Creating new presentation with 16:9 aspect ratio")
            prs = Presentation()
            prs.slide_width = self.SLIDE_WIDTH
            prs.slide_height = self.SLIDE_HEIGHT

            # Add a default blank layout if not already present
            print(f"Available slide layouts: {len(prs.slide_layouts)}")
            blank_layout = prs.slide_layouts[6]  # Use completely blank layout
            print(f"Using layout: {blank_layout}")

            # Override output path with environment variable if available
            if os.getenv("LOCAL_OUTPUT_PATH"):
                print("Overriding output path with LOCAL_OUTPUT_PATH")
                presentation.output_path = os.getenv(
                    "LOCAL_OUTPUT_PATH", "presentation.pptx"
                )
                presentation.output_path = presentation.output_path.replace(".pptx", "")
                print(f"New output path: {presentation.output_path}")

            # Get logo path once for all slides
            print("Fetching logo path...")
            presentation.logo_path = self.slide_handler.get_logo_image_path(
                presentation
            )
            print(f"Logo path: {presentation.logo_path}")

            print(f"Using search source: {presentation.search_source}")
            if presentation.search_source == "youtube":
                print("Processing topic with YouTube agent...")
                extra_content = YouTubeAgent(self.content_generator).process_topic(
                    presentation
                )
            else:
                print("Processing topic with Serper agent...")
                extra_content = SerperAgent(self.content_generator).process_topic(
                    presentation
                )

            # Generate and validate slides
            print("Generating slides content...")
            presentation.slides = self.content_generator.generate_slides(
                presentation, extra_content
            )
            if not presentation.slides:
                raise ValueError("No slides were generated from the configuration")
            print(f"Generated {len(presentation.slides)} slides")

            # Add slides with error handling for each
            for i, slide_data in enumerate(presentation.slides, 1):
                try:
                    print(
                        f"\nProcessing slide {i}/{len(presentation.slides)}: {slide_data.title}"
                    )
                    self.slide_handler.add_slide(
                        prs, blank_layout, slide_data, presentation
                    )
                except Exception as e:
                    print(f"Warning: Error adding slide '{slide_data.title}': {str(e)}")
                    continue

            # Validate presentation before saving
            print("\nValidating presentation...")
            # self.validate_presentation(prs)
            print(f"Saving presentation to: {presentation.output_path}")
            self.storage_handler.save_presentation(prs, presentation.output_path)
            print("Presentation created successfully!")

        except Exception as e:
            print(f"ERROR: Failed to create presentation: {str(e)}")
            raise


def main():
    try:
        # Get config path from environment variable or use default
        config_path = os.getenv("PRESENTATION_CONFIG_PATH", "presentation-config.json")

        # Initialize and run presentation maker
        presentation_maker = PresentationMaker()
        presentation_maker.create_presentation(config_path)

    except Exception as e:
        print(f"Error creating presentation: {str(e)}")
        raise


if __name__ == "__main__":
    main()
